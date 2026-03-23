"""
文档解析服务 - 支持 PDF、Word、TXT 等格式
"""

import io
import os
from typing import Optional, Tuple
from pathlib import Path

import PyPDF2
import pdfplumber
from docx import Document
from pptx import Presentation

from app.core.exceptions import DocumentParseException
import structlog

logger = structlog.get_logger()

class DocumentParser:
    """文档解析器"""

    @staticmethod
    def parse_pdf(file_path: str) -> str:
        """
        解析 PDF 文件，优先使用 pdfplumber 提取文本
        """
        try:
            # 首先尝试 pdfplumber（提取效果好）
            text = DocumentParser._parse_pdf_plumber(file_path)
            if text and len(text.strip()) > 100:
                return text

            # 回退到 PyPDF2
            logger.info("pdfplumber_extracted_insufficient,faking_pypdf2", file=file_path)
            return DocumentParser._parse_pdf_pypdf(file_path)

        except Exception as e:
            logger.error("pdf_parse_failed", file=file_path, error=str(e))
            raise DocumentParseException(
                f"PDF 解析失败: {str(e)}",
                file_type="pdf",
                file_path=file_path
            )

    @staticmethod
    def parse_pdf_plumber(file_path: str) -> str:
        """使用 pdfplumber 解析 PDF"""
        text_parts = []
        with pdfplumber.open(file_path) as pdf:
            for page in pdf.pages:
                page_text = page.extract_text()
                if page_text:
                    text_parts.append(page_text)
        return "\n".join(text_parts)

    @staticmethod
    def _parse_pdf_pypdf(file_path: str) -> str:
        """使用 PyPDF2 解析 PDF"""
        text_parts = []
        with open(file_path, 'rb') as f:
            reader = PyPDF2.PdfReader(f)
            for page in reader.pages:
                text = page.extract_text()
                if text:
                    text_parts.append(text)
        return "\n".join(text_parts)

    @staticmethod
    def parse_docx(file_path: str) -> str:
        """
        解析 Word 文档 (.docx)
        """
        try:
            doc = Document(file_path)
            text_parts = []

            for para in doc.paragraphs:
                if para.text:
                    text_parts.append(para.text)

            # 提取表格内容
            for table in doc.tables:
                for row in table.rows:
                    row_text = " | ".join(cell.text for cell in row.cells)
                    text_parts.append(row_text)

            return "\n".join(text_parts)

        except Exception as e:
            logger.error("docx_parse_failed", file=file_path, error=str(e))
            raise DocumentParseException(
                f"Word 解析失败: {str(e)}",
                file_type="docx",
                file_path=file_path
            )

    @staticmethod
    def parse_pptx(file_path: str) -> str:
        """
        解析 PowerPoint 文档 (.pptx)
        """
        try:
            prs = Presentation(file_path)
            text_parts = []

            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text:
                        text_parts.append(shape.text)

            return "\n".join(text_parts)

        except Exception as e:
            logger.error("pptx_parse_failed", file=file_path, error=str(e))
            raise DocumentParseException(
                f"PPT 解析失败: {str(e)}",
                file_type="pptx",
                file_path=file_path
            )

    @staticmethod
    def parse_txt(file_path: str, encoding: str = "utf-8") -> str:
        """
        解析纯文本文件
        """
        try:
            with open(file_path, 'r', encoding=encoding) as f:
                return f.read()
        except UnicodeDecodeError:
            # 尝试其他编码
            for enc in ["gbk", "gb2312", "latin-1", "utf-8-sig"]:
                try:
                    with open(file_path, 'r', encoding=enc) as f:
                        return f.read()
                    except UnicodeDecodeError:
                        continue
            raise DocumentParseException(
                "无法解析文本文件编码",
                file_type="txt",
                file_path=file_path
                       except Exception as e:
            logger.error("txt_parse_failed", file=file_path, error=str(e))
            raise DocumentParseException(
                f"TXT 解析失败: {str(e)}",
                file_type="txt",
                file_path=file_path
            )

    @staticmethod
    def parse_file(file_path: str) -> Tuple[str, str]:
        """
        自动检测文件类型并解析

        Args:
            file_path: 文件路径

        Returns:
            Tuple[提取的文本内容, 文件类型]
        """
        if not os.path.exists(file_path):
            raise DocumentParseException("文件不存在", file_path=file_path)

        ext = Path(file_path).suffix.lower()

        parsers = {
            '.pdf': DocumentParser.parse_pdf,
            '.docx': DocumentParser.parse_docx,
            '.doc': lambda p: DocumentParser.parse_docx(p),  # TODO: 需 antiword 或其他工具
            '.pptx': DocumentParser.parse_pptx,
            '.txt': DocumentParser.parse_txt,
            '.md': DocumentParser.parse_txt,
        }

        parser = parsers.get(ext)
        if not parser:
            raise DocumentParseException(
                f"不支持的文件类型: {ext}",
                file_type=ext,
                file_path=file_path
            )

        logger.info("parsing_document", file=file_path, type=ext)
        text = parser(file_path)

        # 清理文本
        text = DocumentParser._clean_text(text)

        if not text or len(text.strip()) < 10:
            raise DocumentParseException(
                "提取的文本内容过少或为空",
                file_type=ext,
                file_path=file_path
            )

        return text, ext

    @staticmethod
    def _clean_text(text: str) -> str:
        """清理提取的文本"""
        # 移除过多空白行
        lines = text.split('\n')
        cleaned_lines = []
        blank_count = 0

        for line in lines:
            stripped = line.strip()
            if not stripped:
                blank_count += 1
                if blank_count <= 2:  # 保留最多2个空行
                    cleaned_lines.append("")
            else:
                blank_count = 0
                cleaned_lines.append(stripped)

        # 移除首尾空行
        while cleaned_lines and not cleaned_lines[0]:
            cleaned_lines.pop(0)
        while cleaned_lines and not cleaned_lines[-1]:
            cleaned_lines.pop()

        return "\n".join(cleaned_lines)

    @staticmethod
    def get_supported_extensions() -> List[str]:
        """获取支持的文档格式列表"""
        return ['.pdf', '.docx', '.pptx', '.txt', '.md']
